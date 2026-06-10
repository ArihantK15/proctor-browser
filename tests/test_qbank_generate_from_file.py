"""Tests for AI-questions-from-notes: PPTX extraction + generate-from-file."""
import io
import os
import sys
from unittest.mock import patch, AsyncMock

sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))
os.environ.setdefault("SUPABASE_URL", "https://fake.supabase.co")
os.environ.setdefault("SUPABASE_SERVICE_ROLE_KEY", "fake-key")
os.environ.setdefault("SUPABASE_JWT_SECRET", "test-secret-key-at-least-32-chars-long!!")

from tests.conftest import make_admin_token  # noqa: E402
from app.parsers.document import extract_document, UnreadableDocError  # noqa: E402

TEACHER = {"id": "teacher-1", "email": "p@t.com", "org_id": "org-1",
           "org_role": "teacher", "full_name": "P", "status": "active"}


def _hdr():
    return {"Authorization": f"Bearer {make_admin_token(teacher_id='teacher-1')}"}


def _pptx(slides):
    from pptx import Presentation
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for body, note in slides:
        s = prs.slides.add_slide(blank)
        tb = s.shapes.add_textbox(0, 0, 100, 100).text_frame
        tb.text = body
        if note:
            s.notes_slide.notes_text_frame.text = note
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── PPTX extraction ────────────────────────────────────────────────

def test_pptx_extracts_slide_text_and_notes():
    data = _pptx([("Photosynthesis basics", "Light and dark reactions"),
                  ("Calvin cycle", "")])
    doc = extract_document(data, "deck.pptx")
    assert doc.kind == "pptx"
    assert "Photosynthesis basics" in doc.text
    assert "Light and dark reactions" in doc.text   # speaker note included
    assert "Calvin cycle" in doc.text
    assert doc.pdf_bytes is None


def test_pptx_slide_order_preserved():
    data = _pptx([("First slide", ""), ("Second slide", "")])
    doc = extract_document(data, "deck.pptx")
    assert doc.text.index("First slide") < doc.text.index("Second slide")


def test_unknown_extension_still_rejected():
    import pytest
    with pytest.raises(UnreadableDocError):
        extract_document(b"x", "notes.key")


# ── /generate-from-file endpoint ──────────────────────────────────

_FAKE_QS = [{"question": "What drives photosynthesis?", "question_type": "mcq_single",
             "option_A": "Light", "option_B": "Sound", "option_C": "Heat",
             "option_D": "Sugar", "correct": "A", "tags": ["bio"]}]


def test_requires_auth(client):
    r = client.post("/api/v1/admin/question-bank/generate-from-file",
                    files={"file": ("d.pptx", b"x", "application/vnd.ms-powerpoint")})
    assert r.status_code == 401


def test_generate_from_pptx_returns_preview(client):
    data = _pptx([("Photosynthesis is how plants make energy from light. "
                   "It has light and dark reactions in the chloroplast.", "")])
    with patch("app.auth.admin_auth._get_teacher_by_id", return_value=TEACHER), \
         patch("app.llm.is_configured", return_value=True), \
         patch("app.llm.generate_questions", AsyncMock(return_value=_FAKE_QS)):
        r = client.post("/api/v1/admin/question-bank/generate-from-file",
                        files={"file": ("deck.pptx", data,
                               "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
                        headers=_hdr())
    assert r.status_code == 200, r.text
    body = r.json()
    assert body["count"] == 1
    assert body["questions"][0]["correct"] == "A"
    assert body["truncated"] is False


def test_llm_not_configured_503(client):
    data = _pptx([("Enough text here to pass the minimum length guard easily.", "")])
    with patch("app.auth.admin_auth._get_teacher_by_id", return_value=TEACHER), \
         patch("app.llm.is_configured", return_value=False):
        r = client.post("/api/v1/admin/question-bank/generate-from-file",
                        files={"file": ("deck.pptx", data, "application/octet-stream")},
                        headers=_hdr())
    assert r.status_code == 503


def test_bad_extension_415(client):
    with patch("app.auth.admin_auth._get_teacher_by_id", return_value=TEACHER), \
         patch("app.llm.is_configured", return_value=True):
        r = client.post("/api/v1/admin/question-bank/generate-from-file",
                        files={"file": ("notes.txt", b"hello world", "text/plain")},
                        headers=_hdr())
    assert r.status_code == 415


def test_too_little_text_422(client):
    data = _pptx([("hi", "")])
    with patch("app.auth.admin_auth._get_teacher_by_id", return_value=TEACHER), \
         patch("app.llm.is_configured", return_value=True), \
         patch("app.llm.generate_questions", AsyncMock(return_value=_FAKE_QS)):
        r = client.post("/api/v1/admin/question-bank/generate-from-file",
                        files={"file": ("deck.pptx", data, "application/octet-stream")},
                        headers=_hdr())
    assert r.status_code == 422


# ── source-text cleanup ────────────────────────────────────────────

def test_clean_source_text_strips_ligature_noise_keeps_code():
    from app.routers.question_bank import _clean_source_text
    raw = "Heading\nfi\nState vs Props\nff\nconst n = 5;\n\n\n\nDone"
    out = _clean_source_text(raw)
    assert "\nfi\n" not in out and "\nff\n" not in out   # standalone noise gone
    assert "const n = 5;" in out                          # code variable untouched
    assert "\n\n\n" not in out                            # collapsed blank runs
