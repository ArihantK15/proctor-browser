"""Document → normalized text + line layout. PDF (pdfplumber) and DOCX (python-docx).

Raises ScannedPdfError for image-only PDFs and UnreadableDocError for unknown /
corrupt inputs. Keeps the original PDF bytes for later region rasterization.
Pure of the rest of the app — knows nothing about questions.
"""
import io
import logging
from dataclasses import dataclass, field

logger = logging.getLogger("qbank.document")


class ScannedPdfError(Exception):
    """PDF has effectively no extractable text (scanned/image-only)."""


class UnreadableDocError(Exception):
    """Unknown extension or corrupt/encrypted file."""


@dataclass
class DocLine:
    text: str
    page: int
    bbox: list | None
    fonts: set = field(default_factory=set)
    has_image: bool = False


@dataclass
class ExtractedDoc:
    lines: list
    pdf_bytes: bytes | None
    kind: str

    @property
    def text(self) -> str:
        return "\n".join(l.text for l in self.lines)


_MIN_CHARS_FOR_TEXT_PDF = 20   # total extractable chars below this → scanned-like


def _extract_pdf(data: bytes) -> ExtractedDoc:
    import pdfplumber
    lines: list = []
    total_chars = 0
    try:
        with pdfplumber.open(io.BytesIO(data)) as pdf:
            for pidx, page in enumerate(pdf.pages):
                images = page.images or []
                words = page.extract_words(
                    use_text_flow=True, keep_blank_chars=False,
                    extra_attrs=["fontname"]) or []
                # Cluster words into visual lines by their 'top' coordinate.
                rows: dict = {}
                for w in words:
                    key = round(w["top"] / 3.0)   # ~3pt bucket
                    rows.setdefault(key, []).append(w)
                for key in sorted(rows):
                    ws = sorted(rows[key], key=lambda w: w["x0"])
                    text = " ".join(w["text"] for w in ws)
                    total_chars += len(text)
                    x0 = min(w["x0"] for w in ws)
                    x1 = max(w["x1"] for w in ws)
                    top = min(w["top"] for w in ws)
                    bot = max(w["bottom"] for w in ws)
                    fonts = {w.get("fontname", "") for w in ws}
                    bbox = [x0, top, x1, bot]
                    has_img = any(
                        not (im["x1"] < x0 or im["x0"] > x1 or
                             im["bottom"] < top or im["top"] > bot)
                        for im in images)
                    lines.append(DocLine(text=text, page=pidx, bbox=bbox,
                                         fonts=fonts, has_image=has_img))
    except Exception as e:
        logger.warning("[qbank] PDF open failed: %s", e)
        raise UnreadableDocError("Couldn't open this file.") from e
    if total_chars < _MIN_CHARS_FOR_TEXT_PDF:
        raise ScannedPdfError("This looks like a scanned PDF.")
    return ExtractedDoc(lines=lines, pdf_bytes=data, kind="pdf")


def _extract_docx(data: bytes) -> ExtractedDoc:
    import docx
    try:
        d = docx.Document(io.BytesIO(data))
    except Exception as e:
        raise UnreadableDocError("Couldn't open this file.") from e
    lines = [DocLine(text=p.text, page=0, bbox=None, fonts=set())
             for p in d.paragraphs if p.text and p.text.strip()]
    return ExtractedDoc(lines=lines, pdf_bytes=None, kind="docx")


def _extract_pptx(data: bytes) -> ExtractedDoc:
    from pptx import Presentation
    try:
        prs = Presentation(io.BytesIO(data))
    except Exception as e:
        raise UnreadableDocError("Couldn't open this file.") from e
    lines: list = []
    for sidx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            # Tables: pull each cell, row by row.
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text and c.text.strip()]
                    if cells:
                        lines.append(DocLine(text=" ".join(cells), page=sidx, bbox=None))
                continue
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    txt = "".join(r.text for r in para.runs).strip() or (para.text or "").strip()
                    if txt:
                        lines.append(DocLine(text=txt, page=sidx, bbox=None))
        # Speaker notes — often the richest prose for question generation.
        if slide.has_notes_slide:
            notes = (slide.notes_slide.notes_text_frame.text or "").strip()
            if notes:
                lines.append(DocLine(text=notes, page=sidx, bbox=None))
    return ExtractedDoc(lines=lines, pdf_bytes=None, kind="pptx")


def extract_document(data: bytes, filename: str) -> ExtractedDoc:
    name = (filename or "").lower()
    if name.endswith(".pdf"):
        return _extract_pdf(data)
    if name.endswith(".docx"):
        return _extract_docx(data)
    if name.endswith(".pptx"):
        return _extract_pptx(data)
    raise UnreadableDocError("Only PDF, Word (.docx) and PowerPoint (.pptx) files are supported.")
